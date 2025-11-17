"""
PowerPoint Template Generator for ModelIt K12
Creates professional branded slide templates for conferences and webinars
Output: /assets/templates/*.potx
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Import brand constants
sys.path.append(str(Path(__file__).parent))
from brand_constants import COLORS_PPT, PPT_LAYOUT, FONTS, FONT_SIZES, BRAND_INFO

def create_title_template(prs):
    """Create title slide template"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS_PPT["background_light"]

    # Add colored accent bar at top
    accent_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0),
        Inches(0),
        Inches(PPT_LAYOUT["slide_width"]),
        Inches(0.15)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS_PPT["primary_light_blue"]
    accent_bar.line.width = Pt(0)

    # Title placeholder
    title_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(2.5),
        Inches(8.0),
        Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "[Presentation Title]"
    title_para.font.size = Pt(FONT_SIZES["title"])
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS_PPT["primary_dark_blue"]
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle placeholder
    subtitle_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(4.2),
        Inches(8.0),
        Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "[Subtitle or Speaker Name]"
    subtitle_para.font.size = Pt(FONT_SIZES["heading_2"])
    subtitle_para.font.color.rgb = COLORS_PPT["secondary_navy"]
    subtitle_para.alignment = PP_ALIGN.CENTER

    # ModelIt branding at bottom
    brand_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(6.5),
        Inches(8.0),
        Inches(0.5)
    )
    brand_frame = brand_box.text_frame
    brand_para = brand_frame.paragraphs[0]
    brand_para.text = f"{BRAND_INFO['name']} | {BRAND_INFO['tagline']}"
    brand_para.font.size = Pt(14)
    brand_para.font.italic = True
    brand_para.font.color.rgb = COLORS_PPT["accent_teal"]
    brand_para.alignment = PP_ALIGN.CENTER

def create_content_template(prs):
    """Create standard content slide template"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Top accent bar
    accent_bar = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        Inches(PPT_LAYOUT["slide_width"]),
        Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS_PPT["primary_light_blue"]
    accent_bar.line.width = Pt(0)

    # Slide title
    title_box = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(0.5),
        Inches(8.5),
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "[Slide Title]"
    title_para.font.size = Pt(FONT_SIZES["heading_1"])
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS_PPT["primary_dark_blue"]

    # Content area (bullet points)
    content_box = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(1.8),
        Inches(8.5),
        Inches(4.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Add sample bullets
    bullets = [
        "Key point 1: Replace with your content",
        "Key point 2: Systems thinking approach",
        "Key point 3: NGSS alignment",
        "Key point 4: Cell Collective integration"
    ]

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()

        para.text = bullet_text
        para.level = 0
        para.font.size = Pt(FONT_SIZES["body"])
        para.font.color.rgb = COLORS_PPT["secondary_navy"]

    # Footer with page number
    footer_box = slide.shapes.add_textbox(
        Inches(8.5),
        Inches(7.0),
        Inches(1.0),
        Inches(0.3)
    )
    footer_frame = footer_box.text_frame
    footer_para = footer_frame.paragraphs[0]
    footer_para.text = "#"
    footer_para.font.size = Pt(12)
    footer_para.font.color.rgb = COLORS_PPT["secondary_navy"]
    footer_para.alignment = PP_ALIGN.RIGHT

def create_section_header_template(prs):
    """Create section divider slide template"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background with gradient effect (simulate with two rectangles)
    bg_rect1 = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        Inches(PPT_LAYOUT["slide_width"]),
        Inches(PPT_LAYOUT["slide_height"] / 2)
    )
    bg_rect1.fill.solid()
    bg_rect1.fill.fore_color.rgb = COLORS_PPT["primary_dark_blue"]
    bg_rect1.line.width = Pt(0)

    bg_rect2 = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(PPT_LAYOUT["slide_height"] / 2),
        Inches(PPT_LAYOUT["slide_width"]),
        Inches(PPT_LAYOUT["slide_height"] / 2)
    )
    bg_rect2.fill.solid()
    bg_rect2.fill.fore_color.rgb = COLORS_PPT["primary_light_blue"]
    bg_rect2.line.width = Pt(0)

    # Section title
    section_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(3.0),
        Inches(8.0),
        Inches(1.5)
    )
    section_frame = section_box.text_frame
    section_para = section_frame.paragraphs[0]
    section_para.text = "[Section Title]"
    section_para.font.size = Pt(54)
    section_para.font.bold = True
    section_para.font.color.rgb = RGBColor(255, 255, 255)
    section_para.alignment = PP_ALIGN.CENTER

def create_two_column_template(prs):
    """Create two-column layout template"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Top accent
    accent_bar = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        Inches(PPT_LAYOUT["slide_width"]),
        Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS_PPT["accent_teal"]
    accent_bar.line.width = Pt(0)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(0.5),
        Inches(8.5),
        Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "[Two-Column Layout]"
    title_para.font.size = Pt(FONT_SIZES["heading_1"])
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS_PPT["primary_dark_blue"]

    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(1.8),
        Inches(4.0),
        Inches(4.5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_para = left_frame.paragraphs[0]
    left_para.text = "[Left Column Content]"
    left_para.font.size = Pt(FONT_SIZES["body"])
    left_para.font.color.rgb = COLORS_PPT["secondary_navy"]

    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(5.25),
        Inches(1.8),
        Inches(4.0),
        Inches(4.5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_para = right_frame.paragraphs[0]
    right_para.text = "[Right Column Content]"
    right_para.font.size = Pt(FONT_SIZES["body"])
    right_para.font.color.rgb = COLORS_PPT["secondary_navy"]

    # Vertical divider line
    divider = slide.shapes.add_shape(
        1,  # Rectangle (thin)
        Inches(4.95),
        Inches(1.8),
        Inches(0.05),
        Inches(4.5)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLORS_PPT["background_light"]
    divider.line.width = Pt(0)

def generate_ppt_templates():
    """Generate all PowerPoint templates"""

    print("\n📊 Generating ModelIt K12 PowerPoint Templates...")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(PPT_LAYOUT["slide_width"])
    prs.slide_height = Inches(PPT_LAYOUT["slide_height"])

    # Generate templates
    templates = [
        ("Title Slide", create_title_template),
        ("Content Slide", create_content_template),
        ("Section Header", create_section_header_template),
        ("Two-Column Layout", create_two_column_template)
    ]

    for template_name, template_func in templates:
        print(f"  ├─ Creating {template_name} template...")
        template_func(prs)

    # Save as template
    output_path = Path(__file__).parent.parent / "assets" / "templates" / "modelit_presentation_templates.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"  └─ ✅ Saved to: {output_path}")
    print(f"\n📊 PowerPoint Templates Complete!")
    print(f"   • {len(templates)} template slides")
    print(f"   • Branded colors and fonts")
    print(f"   • Ready for conference/webinar use")

    return output_path

if __name__ == "__main__":
    generate_ppt_templates()
