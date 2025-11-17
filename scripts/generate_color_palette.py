"""
Color Palette Generator for ModelIt K12
Creates professional PowerPoint presentation showing complete brand color system
Output: modelit_color_palette.pptx
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

# Import brand constants
sys.path.append(str(Path(__file__).parent))
from brand_constants import BRAND_COLORS, COLORS_PPT, PPT_LAYOUT, BRAND_INFO

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_color_swatch(slide, color_info, x_position, y_position, width=2.5, height=1.5):
    """Add a color swatch box with details"""

    # Get RGB values
    rgb = color_info["rgb"]

    # Add colored rectangle
    swatch = slide.shapes.add_shape(
        1,  # Rectangle shape
        Inches(x_position),
        Inches(y_position),
        Inches(width),
        Inches(height)
    )
    swatch.fill.solid()
    swatch.fill.fore_color.rgb = RGBColor(*rgb)
    swatch.line.width = Pt(2)
    swatch.line.color.rgb = COLORS_PPT["secondary_navy"]

    # Add color name above swatch
    name_box = slide.shapes.add_textbox(
        Inches(x_position),
        Inches(y_position - 0.4),
        Inches(width),
        Inches(0.35)
    )
    name_frame = name_box.text_frame
    name_frame.text = color_info["name"]
    name_para = name_frame.paragraphs[0]
    name_para.font.size = Pt(16)
    name_para.font.bold = True
    name_para.font.color.rgb = COLORS_PPT["secondary_navy"]
    name_para.alignment = PP_ALIGN.CENTER

    # Add color codes below swatch
    codes_box = slide.shapes.add_textbox(
        Inches(x_position),
        Inches(y_position + height + 0.1),
        Inches(width),
        Inches(0.8)
    )
    codes_frame = codes_box.text_frame
    codes_frame.word_wrap = True

    # HEX code
    p1 = codes_frame.paragraphs[0]
    p1.text = f"HEX: {color_info['hex']}"
    p1.font.size = Pt(11)
    p1.font.name = "Consolas"
    p1.alignment = PP_ALIGN.CENTER

    # RGB code
    p2 = codes_frame.add_paragraph()
    p2.text = f"RGB: {rgb[0]}, {rgb[1]}, {rgb[2]}"
    p2.font.size = Pt(11)
    p2.font.name = "Consolas"
    p2.alignment = PP_ALIGN.CENTER

    # CMYK code
    cmyk = color_info['cmyk']
    p3 = codes_frame.add_paragraph()
    p3.text = f"CMYK: {cmyk[0]}, {cmyk[1]}, {cmyk[2]}, {cmyk[3]}"
    p3.font.size = Pt(11)
    p3.font.name = "Consolas"
    p3.alignment = PP_ALIGN.CENTER

    # Usage note
    usage_box = slide.shapes.add_textbox(
        Inches(x_position),
        Inches(y_position + height + 1.0),
        Inches(width),
        Inches(0.6)
    )
    usage_frame = usage_box.text_frame
    usage_frame.word_wrap = True
    usage_para = usage_frame.paragraphs[0]
    usage_para.text = color_info["usage"]
    usage_para.font.size = Pt(10)
    usage_para.font.italic = True
    usage_para.font.color.rgb = COLORS_PPT["secondary_navy"]
    usage_para.alignment = PP_ALIGN.CENTER

def create_title_slide(prs):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS_PPT["background_light"]

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(PPT_LAYOUT["margin_left"]),
        Inches(2.0),
        Inches(PPT_LAYOUT["slide_width"] - PPT_LAYOUT["margin_left"] - PPT_LAYOUT["margin_right"]),
        Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "ModelIt! K12"
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS_PPT["primary_dark_blue"]
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(PPT_LAYOUT["margin_left"]),
        Inches(3.7),
        Inches(PPT_LAYOUT["slide_width"] - PPT_LAYOUT["margin_left"] - PPT_LAYOUT["margin_right"]),
        Inches(1.0)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "Brand Color Palette"
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = COLORS_PPT["primary_light_blue"]
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide.shapes.add_textbox(
        Inches(PPT_LAYOUT["margin_left"]),
        Inches(5.0),
        Inches(PPT_LAYOUT["slide_width"] - PPT_LAYOUT["margin_left"] - PPT_LAYOUT["margin_right"]),
        Inches(0.6)
    )
    tagline_frame = tagline_box.text_frame
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.text = BRAND_INFO["tagline"]
    tagline_para.font.size = Pt(18)
    tagline_para.font.italic = True
    tagline_para.font.color.rgb = COLORS_PPT["secondary_navy"]
    tagline_para.alignment = PP_ALIGN.CENTER

    # Generated date
    date_box = slide.shapes.add_textbox(
        Inches(PPT_LAYOUT["margin_left"]),
        Inches(6.5),
        Inches(PPT_LAYOUT["slide_width"] - PPT_LAYOUT["margin_left"] - PPT_LAYOUT["margin_right"]),
        Inches(0.4)
    )
    date_frame = date_box.text_frame
    date_para = date_frame.paragraphs[0]
    date_para.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
    date_para.font.size = Pt(12)
    date_para.font.color.rgb = COLORS_PPT["secondary_navy"]
    date_para.alignment = PP_ALIGN.CENTER

def create_primary_colors_slide(prs):
    """Create slide showing primary brand colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(PPT_LAYOUT["margin_left"]),
        Inches(0.5),
        Inches(PPT_LAYOUT["slide_width"] - PPT_LAYOUT["margin_left"] - PPT_LAYOUT["margin_right"]),
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Primary Colors"
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS_PPT["primary_dark_blue"]
    title_para.alignment = PP_ALIGN.CENTER

    # Add primary color swatches
    primary_colors = ["primary_dark_blue", "primary_light_blue", "secondary_navy"]
    x_positions = [1.0, 3.75, 6.5]

    for i, color_key in enumerate(primary_colors):
        add_color_swatch(slide, BRAND_COLORS[color_key], x_positions[i], 2.0)

def create_accent_colors_slide(prs):
    """Create slide showing accent and background colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(PPT_LAYOUT["margin_left"]),
        Inches(0.5),
        Inches(PPT_LAYOUT["slide_width"] - PPT_LAYOUT["margin_left"] - PPT_LAYOUT["margin_right"]),
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Accent & Background Colors"
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS_PPT["primary_dark_blue"]
    title_para.alignment = PP_ALIGN.CENTER

    # Add accent color swatches
    accent_colors = ["background_light", "accent_teal", "accent_gold"]
    x_positions = [1.0, 3.75, 6.5]

    for i, color_key in enumerate(accent_colors):
        add_color_swatch(slide, BRAND_COLORS[color_key], x_positions[i], 2.0)

def create_usage_examples_slide(prs):
    """Create slide showing color usage examples"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS_PPT["background_light"]

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(0.5),
        Inches(8.5),
        Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Color Usage Examples"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS_PPT["primary_dark_blue"]
    title_para.alignment = PP_ALIGN.CENTER

    # Example 1: Heading + Body
    ex1_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(8.0), Inches(1.5))
    ex1_frame = ex1_box.text_frame
    ex1_frame.word_wrap = True

    ex1_heading = ex1_frame.paragraphs[0]
    ex1_heading.text = "Example Heading (Primary Dark Blue)"
    ex1_heading.font.size = Pt(28)
    ex1_heading.font.bold = True
    ex1_heading.font.color.rgb = COLORS_PPT["primary_dark_blue"]

    ex1_body = ex1_frame.add_paragraph()
    ex1_body.text = "This is body text in the same primary dark blue. Use for main content and professional materials."
    ex1_body.font.size = Pt(16)
    ex1_body.font.color.rgb = COLORS_PPT["secondary_navy"]

    # Example 2: Call to Action
    ex2_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1.5),
        Inches(3.8),
        Inches(3.0),
        Inches(0.7)
    )
    ex2_box.fill.solid()
    ex2_box.fill.fore_color.rgb = COLORS_PPT["primary_light_blue"]
    ex2_box.line.width = Pt(0)

    ex2_text = ex2_box.text_frame
    ex2_para = ex2_text.paragraphs[0]
    ex2_para.text = "Download Now"
    ex2_para.font.size = Pt(24)
    ex2_para.font.bold = True
    ex2_para.font.color.rgb = RGBColor(255, 255, 255)
    ex2_para.alignment = PP_ALIGN.CENTER

    # Example 3: Highlight/Warning
    ex3_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(5.5),
        Inches(3.8),
        Inches(3.0),
        Inches(0.7)
    )
    ex3_box.fill.solid()
    ex3_box.fill.fore_color.rgb = COLORS_PPT["accent_gold"]
    ex3_box.line.width = Pt(0)

    ex3_text = ex3_box.text_frame
    ex3_para = ex3_text.paragraphs[0]
    ex3_para.text = "⚠️ Important Note"
    ex3_para.font.size = Pt(22)
    ex3_para.font.bold = True
    ex3_para.font.color.rgb = COLORS_PPT["secondary_navy"]
    ex3_para.alignment = PP_ALIGN.CENTER

    # Example 4: Science Theme
    ex4_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(8.0), Inches(1.5))
    ex4_frame = ex4_box.text_frame
    ex4_frame.word_wrap = True

    ex4_heading = ex4_frame.paragraphs[0]
    ex4_heading.text = "🔬 Science Content (Accent Teal)"
    ex4_heading.font.size = Pt(24)
    ex4_heading.font.bold = True
    ex4_heading.font.color.rgb = COLORS_PPT["accent_teal"]

    ex4_body = ex4_frame.add_paragraph()
    ex4_body.text = "Use accent teal for scientific content, molecular diagrams, and biology-themed materials."
    ex4_body.font.size = Pt(16)
    ex4_body.font.color.rgb = COLORS_PPT["secondary_navy"]

def generate_color_palette():
    """Main function to generate color palette PowerPoint"""

    print("🎨 Generating ModelIt K12 Color Palette...")

    # Create presentation with widescreen layout
    prs = Presentation()
    prs.slide_width = Inches(PPT_LAYOUT["slide_width"])
    prs.slide_height = Inches(PPT_LAYOUT["slide_height"])

    # Add slides
    print("  ├─ Creating title slide...")
    create_title_slide(prs)

    print("  ├─ Creating primary colors slide...")
    create_primary_colors_slide(prs)

    print("  ├─ Creating accent colors slide...")
    create_accent_colors_slide(prs)

    print("  ├─ Creating usage examples slide...")
    create_usage_examples_slide(prs)

    # Save presentation
    output_path = Path(__file__).parent.parent / "assets" / "colors" / "modelit_color_palette.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"  └─ ✅ Saved to: {output_path}")
    print(f"\n📊 Color Palette Complete!")
    print(f"   • 6 brand colors defined")
    print(f"   • 4 slides generated")
    print(f"   • HEX, RGB, and CMYK codes included")
    print(f"   • Usage examples provided")

    return output_path

if __name__ == "__main__":
    generate_color_palette()
